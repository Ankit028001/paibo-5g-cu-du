#!/usr/bin/env python3
"""
generate_paibo_pptx.py

READ-ONLY against all existing data. Builds a single PPTX presenting:
  - REAL OAI CU-DU measured results (from /opt/oai/openairinterface5g/phase2/)
  - ns-3 / 5G-LENA discrete-event SIMULATION results (from ns3_cudu_phase/)
  - PAIBO Result Types 1-4 (NOT AVAILABLE -- no PAIBO pipeline exists on
    this machine; nothing here is fabricated)

IMPORTANT FACT ESTABLISHED BEFORE WRITING THIS SCRIPT (do not re-derive
differently later): no "16-UE OAI measured" dataset exists anywhere on
this machine. The only OAI CU-DU experiment with a structured per-UE CSV
table is the 8-UE / 106-PRB pilot
(phase2/20260902_vrtsim_cudu_8ue_106prb/), whose real, measured result was
0/8 UEs attached (PHY-layer sync failure). Separately, in a different
100-UE-configured experiment (phase2/20260903_100ue_vrtsim_cudu_189prb_patched/),
the peak real OAI CU-DU registration achieved on this machine was 2 UEs
(registered + PDU session established) -- narrated in that run's
SUMMARY.md, not captured in a structured per-UE CSV with radio KPIs.
Both facts are presented explicitly in the deck; neither is glossed over.

No ns-3 or OAI simulation is run by this script. No source file is
modified. Nothing is pushed to GitHub.
"""

import csv
import os
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = "/mnt/c/Users/Common/Downloads/Siya"
OAI_8UE = "/opt/oai/openairinterface5g/phase2/20260902_vrtsim_cudu_8ue_106prb"
CUDU = os.path.join(BASE, "ns3_cudu_phase")

DARK_GRAY = RGBColor(0x3B, 0x3B, 0x3B)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BORDER_GRAY = RGBColor(0x88, 0x88, 0x88)
TEXT_DARK = RGBColor(0x20, 0x20, 0x20)
RED_FLAG = RGBColor(0xB0, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    # green accent bar under title
    bar = slide.shapes.add_shape(1, Inches(0.4), Inches(0.85), Inches(3.0), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.4))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.italic = True
        sp.font.color.rgb = RGBColor(0x60, 0x60, 0x60)


def add_source_label(slide, text, flagged=False, top=Inches(7.0)):
    box = slide.shapes.add_textbox(Inches(0.4), top, Inches(12.5), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = RED_FLAG if flagged else RGBColor(0x50, 0x50, 0x50)


def add_notes(slide, lines, top, width=Inches(12.5), size=11):
    box = slide.shapes.add_textbox(Inches(0.4), top, width, Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(0x40, 0x40, 0x40)


def add_table(slide, rows, left, top, width, height, col_widths=None, header_row=True,
              font_size=11, na_color=RGBColor(0x99, 0x33, 0x00)):
    n_rows = len(rows)
    n_cols = len(rows[0])
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gtable.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(int(width * (w / total)))
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = "" if val is None else str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(font_size if not (header_row and r == 0) else font_size)
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            if header_row and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_GRAY
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
                if str(val).strip().upper() in ("N/A", "NA", "NOT AVAILABLE", "NOT DIRECTLY COMPARABLE",
                                                 "NOT YET MEASURED",
                                                 "NOT AVAILABLE — PAIBO BIP PIPELINE NOT IMPLEMENTED/EXECUTED"):
                    para.font.color.rgb = na_color
                    para.font.italic = True
                else:
                    para.font.color.rgb = TEXT_DARK
    return gtable


# ---------------------------------------------------------------------------
# Data loading (real files only)
# ---------------------------------------------------------------------------

def load_csv(path):
    with open(path, newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def load_cudu_ladder():
    levels = [1, 10, 25, 50, 100, 150, 200]
    wallclock = {1: 5.30673, 10: 25.1937, 25: 78.9827, 50: 191.394,
                 100: 536.604, 150: 1129.12, 200: 1530.57}
    bler = {1: 0.0, 10: 0.0, 25: 0.0, 50: 0.0, 100: 0.0, 150: 0.0, 200: 0.0}
    harq = {1: 0, 10: 0, 25: 0, 50: 0, 100: 0, 150: 0, 200: 0}
    jitter = {1: 0.111348, 10: 0.172503, 25: 0.188590, 50: 0.209987,
              100: 0.258176, 150: 0.307347, 200: 0.371013}
    out = []
    for n in levels:
        row = load_csv(os.path.join(CUDU, f"ue_{n}", "per_cell_kpis.csv"))[0]
        agg_mbps = (int(row["totalRxBytes"]) * 8 / 1e6) / 30.0
        loss_pct = (int(row["totalLostPackets"]) / max(1, int(row["totalRxBytes"]))) * 100.0
        out.append({
            "n": n,
            "connected": f"{row['rrcConnectedCount']}/{row['configuredUeCount']}",
            "agg_mbps": round(agg_mbps, 3),
            "sinr": round(float(row["avgUeSinrDb"]), 2),
            "bler": bler[n],
            "mcs": round(float(row["avgUeMcs"]), 2) if row["avgUeMcs"] else "NA (fullTraces=false)",
            "harq": harq[n],
            "loss_pct": round(loss_pct, 4),
            "jitter_ms": round(jitter[n], 4),
            "runtime_s": wallclock[n],
        })
    return out


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------

def main():
    prs = new_deck()

    # --- Slide: Title ---
    s = blank_slide(prs)
    box = s.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PAIBO Baseline KPI Comparison"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p2 = tf.add_paragraph()
    p2.text = "Real OAI CU-DU Measurements  vs.  ns-3/5G-LENA Discrete-Event Simulation  vs.  PAIBO Requirements"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x50, 0x50, 0x50)
    p3 = tf.add_paragraph()
    p3.text = "No fabricated, extrapolated, or estimated KPI values. Every table states its data source."
    p3.font.size = Pt(13)
    p3.font.italic = True
    p3.font.color.rgb = GREEN
    bar = s.shapes.add_shape(1, Inches(1), Inches(2.15), Inches(4.0), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    # --- Slide: Per-UE KPI Summary (REAL OAI) ---
    s = blank_slide(prs)
    add_title(s, "Per-UE KPI Summary — OAI CU-DU + Core",
              "REAL OAI MEASUREMENT — 8-UE / 106-PRB pilot (20260902_vrtsim_cudu_8ue_106prb)")
    attach = load_csv(os.path.join(OAI_8UE, "attach_results.csv"))
    header = ["UE", "IMSI", "RRC State", "5GMM State", "PDU Session", "Sync Result", "IP", "Notes"]
    rows = [header]
    for r in attach:
        rows.append([r["ue_id"], r["imsi"], r["rrc_state"], r["fivegmm_state"],
                     r["pdu_session"], r["sync_result"], r["ip_address"], r["notes"][:38]])
    add_table(s, rows, Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.2),
              col_widths=[4, 11, 8, 8, 7, 7, 6, 20], font_size=10.5)
    add_notes(s, [
        "SINR / RSRP / BLER / PRB / Throughput / Loss / Jitter / RRC Errors / Bearer Latency: N/A for all 8 UEs — "
        "root cause: 0/8 UEs completed PHY-layer initial sync (\"synch Failed\" retry loop); none reached "
        "RRC_CONNECTED, 5GMM REGISTERED, or PDU session, so no radio-layer KPI was ever generated in this run.",
        "Separately, in a different 100-UE-configured OAI experiment (20260903_100ue_vrtsim_cudu_189prb_patched, "
        "attempt_03), 2 UEs (UE0, UE1) did register and establish a PDU session — the highest real OAI CU-DU "
        "registration count found on this machine. No structured per-UE radio-KPI file exists for those 2 UEs "
        "either (narrated in that run's SUMMARY.md only). No 16-UE OAI measured dataset exists on this machine.",
    ], Inches(4.75), size=11)
    add_source_label(s, "Source: OAI CU-DU logs — REAL MEASUREMENT (attach_results.csv, per_ue_kpis.csv)")

    # --- Slide: Per-Use-Case Traffic Volume (REAL OAI) ---
    s = blank_slide(prs)
    add_title(s, "Per-Use-Case Traffic Volume — OAI Measured",
              "REAL OAI MEASUREMENT — 8-UE / 106-PRB pilot")
    traffic = load_csv(os.path.join(OAI_8UE, "traffic_results.csv"))
    status_row = next(r for r in traffic if r["metric"] == "status")
    classes = ["mMTC", "Web application", "Mobile application", "Video on Demand", "Live Video", "V2X", "TOTAL"]
    header = ["Use Case", "UEs", "Measured Bytes", "Traffic % (measured)", "Target %", "Deviation (pp)"]
    rows = [header]
    for cls in classes:
        rows.append([cls, "N/A", "N/A", "N/A", "N/A", "N/A"])
    add_table(s, rows, Inches(0.4), Inches(1.4), Inches(9.0), Inches(3.2), font_size=12)
    add_notes(s, [
        f"OAI traffic test status: {status_row['value']}. {status_row['note']}",
        "No OAI measured traffic-volume-by-class data exists on this machine — the control-plane gate "
        "(successful attach) was never passed in this experiment, so no application traffic test ran.",
        "Measured ns-3 traffic-class byte-volume percentages exist separately — see the "
        "\"Traffic Model — OAI vs ns-3 Baseline\" slide.",
    ], Inches(4.9), size=12)
    add_source_label(s, "Source: OAI CU-DU logs — REAL MEASUREMENT (traffic_results.csv) — status: NOT_RUN")

    # --- Slide: Cell-Level KPIs (REAL OAI) ---
    s = blank_slide(prs)
    add_title(s, "Cell-Level KPIs — OAI CU-DU Run (8-UE / 106-PRB pilot)",
              "REAL OAI MEASUREMENT")
    percell = load_csv(os.path.join(OAI_8UE, "per_cell_kpis.csv"))
    pc = {r["metric"]: r["value"] for r in percell}
    header = ["KPI", "Value", "Notes"]
    rows = [header,
            ["UEs Registered", f"{pc['connected_ue_count']}/8", "0/8 attach success (PHY sync failure, all 8 UEs)"],
            ["Aggregate Throughput", "N/A", "No UE reached PDU session — no data-plane traffic possible"],
            ["Aggregate Packet Loss", "N/A", "No data-plane traffic ever established"],
            ["DU CPU", "172% (multi-core)", "resource_during.csv snapshot during 8-UE sync-retry loop"],
            ["CU CPU", "low single-digit %", "resource_during.csv — CU never under stress in this run"],
            ["DU Memory (RSS)", "~5.11 GiB", "resource_during.csv, DU process, initial launch snapshot"],
            ["CU Memory (RSS)", "~155 MiB", "resource_during.csv, CU process snapshot"],
            ["UE process Memory (RSS)", "~703-720 MiB/UE", "sync-retry loop, NOT steady state (no UE reached connected state)"],
            ["F1 Interface (CU-DU)", "UP", "F1 Setup Request/Response confirmed both directions; SCTP assoc. established"],
            ["F1-U (user plane)", "PASS at GTP-U bind only", "Never exercised end-to-end — no UE ever reached PDU session"],
            ["NGAP (RAN-AMF)", "PASS", "NGSetupRequest/Response succeeded; gNB registered with AMF"],
            ["PFCP (SMF-UPF)", "N/A", "Never triggered — requires a UE PDU session request, none occurred"],
            ["PRB Utilization", "N/A", "No data-plane traffic; PRB utilization requires active scheduling"],
            ]
    add_table(s, rows, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.3),
              col_widths=[7, 6, 20], font_size=11.5)
    add_source_label(s, "Source: OAI CU-DU logs — REAL MEASUREMENT (per_cell_kpis.csv, resource_during.csv, SUMMARY.md)",
                      top=Inches(6.95))

    # --- Slide: KPI Availability Matrix (OAI vs PAIBO) ---
    s = blank_slide(prs)
    add_title(s, "KPI Availability — OAI 5G SA Tested vs PAIBO Patent Requirements",
              "REAL OAI MEASUREMENT — status reflects the actual outcome of OAI testing on this machine")
    kpi_rows = [
        ("PHY / RADIO", None),
        (1, "DL SINR", "PHY trace (if UE connected)", "Not Available", "N/A", "per-TB", "Needed for BIP mobility/QoS features"),
        (2, "UL SINR", "PHY trace (if UE connected)", "Not Available", "N/A", "per-TB", "Same"),
        (3, "DL RSSI/RSRP", "PHY trace (if UE connected)", "Not Available", "N/A", "per-UE", "Needed for mobility trajectory input"),
        (4, "UL RSSI/RSRP", "PHY trace (if UE connected)", "Not Available", "N/A", "per-UE", "Same"),
        ("MAC", None),
        (5, "DL BLER", "MAC trace (if UE connected)", "Not Available", "N/A", "per-TB", "AMC/adaptation context"),
        (6, "UL BLER", "MAC trace (if UE connected)", "Not Available", "N/A", "per-TB", "Same"),
        (7, "MCS (DL)", "MAC scheduler log", "Not Available", "N/A", "per-grant", "AMC/adaptation context"),
        (8, "HARQ Errors", "MAC scheduler log", "Not Available", "N/A", "per-grant", "Adaptation/reliability context"),
        (9, "PRB Utilization", "MAC scheduler log", "Not Available", "N/A", "per-cell", "Resource-efficiency claims"),
        ("RRC / CONTROL PLANE", None),
        (10, "RRC State", "RRC log", "Available", "RRC_IDLE (8/8 UEs)", "per-UE", "Direct input to Bearer Intent Predictor state"),
        (11, "Bearer Setup Latency", "CU/DU/AMF log timestamps", "Not Available", "N/A", "per-event", "Core PAIBO Result Type 1 KPI"),
        (12, "F1 Setup Time", "CU/DU log event", "Partial", "PASS (event only, no elapsed-time number logged)", "per-setup", "CU-DU split validation context"),
        (13, "PDU Session Time", "AMF/SMF log timestamps", "Not Available", "N/A (0/8); event-only for the separate 2-UE case", "per-session", "Core PAIBO context"),
        ("APPLICATION / E2E", None),
        (14, "E2E Latency (RTT)", "ping/iperf3", "Not Available", "N/A", "per-packet", "QoS/SLA validation"),
        (15, "Throughput", "iperf3/traffic test", "Not Available", "N/A (test NOT_RUN)", "per-flow", "QoS/SLA validation"),
        (16, "Packet Loss", "iperf3/traffic test", "Not Available", "N/A (test NOT_RUN)", "per-flow", "QoS/SLA validation"),
        (17, "Jitter", "iperf3/traffic test", "Not Available", "N/A (test NOT_RUN)", "per-flow", "QoS/SLA validation"),
    ]
    header = ["#", "KPI Name", "OAI Source", "Status", "Value Observed", "Granularity", "PAIBO Relevance"]
    rows = [header]
    for item in kpi_rows:
        if item[1] is None:
            rows.append([item[0], "", "", "", "", "", ""])
        else:
            rows.append(list(item))
    tbl_shape = add_table(s, rows, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.6),
              col_widths=[3, 16, 15, 10, 22, 8, 22], font_size=9.5)
    # bold + shade section header rows
    table = tbl_shape.table
    for r_idx, item in enumerate(kpi_rows, start=1):
        if item[1] is None:
            cell = table.cell(r_idx, 0)
            cell.merge(table.cell(r_idx, 6))
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xD9, 0xE8, 0xD9)
            p0 = cell.text_frame.paragraphs[0]
            p0.text = item[0]
            p0.font.bold = True
            p0.font.size = Pt(10.5)
            p0.font.color.rgb = GREEN
    add_source_label(s, "Source: OAI CU-DU logs — REAL MEASUREMENT (this machine's actual OAI test outcome, not a hypothetical capability list)",
                      top=Inches(7.05))

    # --- Slide: ns-3 Baseline ---
    s = blank_slide(prs)
    add_title(s, "ns-3 / 5G-LENA Baseline — CU-DU Scaling",
              "SIMULATION DATA — NOT REAL OAI EXECUTION")
    ladder = load_cudu_ladder()
    header = ["UEs", "RRC-Connected", "Agg DL Tput (Mbps)", "Mean DL SINR (dB)", "Mean DL BLER (%)",
              "Mean DL MCS", "HARQ Retx", "Packet Loss (%)", "Mean Jitter (ms)", "Runtime (s)"]
    rows = [header]
    for lv in ladder:
        rows.append([lv["n"], lv["connected"], lv["agg_mbps"], lv["sinr"], lv["bler"],
                     lv["mcs"], lv["harq"], lv["loss_pct"], lv["jitter_ms"], lv["runtime_s"]])
    add_table(s, rows, Inches(0.4), Inches(1.5), Inches(12.5), Inches(3.4), font_size=11.5)
    add_notes(s, [
        "Ladder complete at 1/10/25/50/100/150/200 UEs; no extrapolation beyond 200 UEs performed or implied.",
        "Mean DL MCS is NA at 200 UEs because that level ran with --fullTraces=false (MAC scheduler trace not "
        "enabled at that level to control I/O volume) — not a measurement failure, a documented run configuration.",
        "IDEAL channel (LOS-only pathloss, no fading/shadowing) drives the near-constant ~63.8 dB SINR and 0% "
        "BLER/HARQ at every level — this is a direct consequence of the channel model, not evidence of general "
        "scheduler robustness under realistic error conditions.",
    ], Inches(5.1), size=11.5)
    add_source_label(s, "Source: 5G-LENA/ns-3 — DISCRETE-EVENT SIMULATION (ns3_cudu_phase/, validated ladder)")

    # --- Slide: OAI vs ns-3 comparison ---
    s = blank_slide(prs)
    add_title(s, "Baseline Comparison — Real OAI vs ns-3/5G-LENA",
              "Distinct measurement contexts — not merged or averaged")
    cmp_rows = [
        ["UE registration", "0/8 (8-UE pilot); 2/100 (separate patched run)", "100% at every level (1-200 UEs)", "Not directly comparable — different scale, architecture, and real-time constraint"],
        ["Aggregate DL throughput", "N/A — no data-plane ever established", f"{ladder[-1]['agg_mbps']} Mbps @ 200 UEs", "Not directly comparable — no OAI data-plane baseline exists"],
        ["DL SINR", "N/A — no UE connected", f"~{ladder[-1]['sinr']} dB (IDEAL channel)", "Not directly comparable"],
        ["DL BLER", "N/A", "0.0% at every level", "Not directly comparable"],
        ["MCS", "N/A", "~27.7-27.9 (near-max, IDEAL channel)", "Not directly comparable"],
        ["HARQ", "N/A", "0 retransmissions at every level", "Not directly comparable"],
        ["Packet loss", "N/A — traffic test NOT_RUN", "0.0% at every level", "Not directly comparable"],
        ["Jitter", "N/A", "0.11-0.37 ms across the ladder", "Not directly comparable"],
        ["Bearer setup latency", "Not logged as a number even for the 2-UE success case", "Not implemented in ns-3", "Not available on either side"],
        ["F1 health", "UP — real F1AP Setup Request/Response, SCTP association", "UP — topological p2p link, carries a heartbeat only, no F1AP", "Not directly comparable — real protocol vs. simulated link"],
        ["NGAP health", "PASS — real NGSetupResponse from stock AMF", "N/A — ns-3's LTE-EPC-style core has no NGAP concept", "Not directly comparable"],
        ["PFCP health", "N/A — never triggered (no PDU session)", "N/A — ns-3 core has no PFCP/5GC concept", "Not available on either side"],
        ["Memory (RSS)", "Measured: DU ~5.1 GiB, ~703-720 MiB/UE", "Not systematically profiled in this study", "Not directly comparable — no ns-3 measurement exists"],
        ["CPU", "Measured: DU climbed 35%→70.6%, real-time bottleneck", "Wall-clock runtime only — no real-time constraint applies", "Not directly comparable — different metric meaning entirely"],
        ["PRB utilization", "N/A — no data-plane traffic", "N/A — not a native 5G-LENA trace field", "Not available on either side"],
    ]
    header = ["KPI", "Real OAI CU-DU", "ns-3 CU-DU", "Status / Interpretation"]
    rows = [header] + cmp_rows
    add_table(s, rows, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.9),
              col_widths=[10, 22, 22, 26], font_size=9.5)
    add_source_label(s, "Sources: OAI CU-DU logs (REAL MEASUREMENT) + 5G-LENA/ns-3 (DISCRETE-EVENT SIMULATION) — never merged or averaged",
                      top=Inches(7.25))

    # --- Slide: PAIBO Result Types ---
    s = blank_slide(prs)
    add_title(s, "PAIBO Result Types — Baseline vs PAIBO",
              "PAIBO pipeline (1_generate_traffic.py / 2_train_bip.py / 3_train_rl_sdap.py / paibo_scaling.cc) "
              "confirmed absent on this machine — nothing below is fabricated")
    paibo_rows = [
        ["Bearer Setup Latency", "100-200 ms (reference figure quoted in PAIBO_Patent_Vartika.pptx, slide 23/31 — "
                                  "NOT independently measured on this OAI/ns-3 setup)",
         "Not available — PAIBO BIP pipeline not implemented/executed",
         "PAIBO patent deck (reference only)", "Not yet measured"],
        ["DRB Count Reduction", "5 QoS flows -> 5 DRBs (reference figure, PAIBO deck slide 30 — NOT measured here)",
         "Not available — PAIBO BIP pipeline not implemented/executed",
         "PAIBO patent deck (reference only)", "Not yet measured"],
        ["MAC-CE Adaptation Latency", "~100 ms full RRCReconfiguration (reference figure, PAIBO deck slide 26 — NOT measured here)",
         "Not available — PAIBO BIP pipeline not implemented/executed",
         "PAIBO patent deck (reference only)", "Not yet measured"],
        ["BIP ML Accuracy", "N/A — no baseline classifier exists to compare against",
         "Not available — PAIBO BIP pipeline not implemented/executed",
         "N/A", "Not yet measured"],
    ]
    header = ["Result Type", "Baseline", "PAIBO", "Source", "Measurement Status"]
    rows = [header] + paibo_rows
    add_table(s, rows, Inches(0.3), Inches(1.6), Inches(12.7), Inches(3.6),
              col_widths=[13, 30, 22, 15, 12], font_size=10.5)
    add_notes(s, [
        "BIP accuracy is explicitly required by the PAIBO assignment/PPT material but cannot be produced without "
        "fabrication: no trained model, no labeled bearer-demand dataset, and no scikit-learn/TensorFlow environment "
        "exists on this machine for this purpose. ns-3 KPI values are never substituted for BIP accuracy.",
    ], Inches(5.5), size=12)
    add_source_label(s, "Source: PAIBO pipeline — NOT AVAILABLE (no scripts found on this machine; baseline figures are PAIBO-deck reference values, not local measurements)",
                      top=Inches(7.1))

    # --- Slide: Traffic Model comparison ---
    s = blank_slide(prs)
    add_title(s, "Traffic Model — OAI vs ns-3 Baseline")
    ns3_measured_150 = {"mMTC": 0.74, "Web": 10.4, "Mobile": 12.8, "VoD": 43.4, "Live": 30.2, "V2X": 2.46}
    target = {"mMTC": 1.0, "Web": 8.0, "Mobile": 10.0, "VoD": 35.0, "Live": 25.0, "V2X": 2.0}
    header = ["Use Case", "OAI Measured Traffic %", "ns-3 Measured Traffic % (N=150)", "Target %", "Difference / Note"]
    rows = [header]
    for cls in ["mMTC", "Web", "Mobile", "VoD", "Live", "V2X"]:
        diff = round(ns3_measured_150[cls] - target[cls], 2)
        rows.append([cls, "N/A (traffic test NOT_RUN)", ns3_measured_150[cls], target[cls], f"{diff:+.2f} pp (ns-3 vs target)"])
    add_table(s, rows, Inches(0.4), Inches(1.5), Inches(11.5), Inches(3.0), font_size=13)
    add_notes(s, [
        "Target % is the configured byte-volume share from the assignment instructions, normalised to 100%.",
        "ns-3 % is a measured byte-volume share (FlowMonitor rxBytes per traffic class, N=150 UE level) — not "
        "computed from the configured UE-count shares (40/15/15/12/13/5%), which are a different axis entirely.",
        "OAI % is N/A because the OAI traffic test never ran (gated behind control-plane attach success, which "
        "was not achieved in the measured run).",
    ], Inches(4.7), size=12)
    add_source_label(s, "Sources: OAI CU-DU logs (REAL MEASUREMENT, status NOT_RUN) + 5G-LENA/ns-3 (DISCRETE-EVENT SIMULATION, measured)")

    # --- Slide: Memory/Scaling ---
    s = blank_slide(prs)
    add_title(s, "System Scaling and Memory Budget")
    header = ["Metric", "OAI (REAL MEASUREMENT)", "ns-3 (SIMULATION)"]
    rows = [header,
            ["Host RAM", "82 GiB", "82 GiB (same host)"],
            ["Peak RAM used", "~20.5 GiB (~25%) with 15 real UE processes (separate 100-UE-configured run, attempt_03)",
             "Not systematically profiled in this study — not claimed"],
            ["Per-UE process RSS", "~872 MiB/UE (isolated 1-UE full-PHY measurement); ~703-720 MiB/UE (8-UE pilot, sync-retry, not steady state)",
             "Not measured — not claimed"],
            ["Observed bottleneck", "CPU / real-time scheduling (DU per-callback O(N) cost in vrtsim), NOT RAM",
             "Wall-clock discrete-event runtime (super-linear with UE count); no real-time constraint applies"],
            ["DU CPU trend", "35% -> 70.6%, still rising when halted (15-UE attempt)", "N/A — not a real-time process"],
            ["CU CPU", "~9%, never stressed", "N/A"],
            ]
    add_table(s, rows, Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.6), col_widths=[10, 30, 30], font_size=11.5)
    add_notes(s, [
        "No AI-pipeline (data dump / pre-processing / inferencing) memory footprint has been measured for this "
        "project, so no revised \"max UE capacity minus AI headroom\" figure is presented — see MEMORY_BUDGET_NOTE.md.",
        "OAI's real bottleneck (CPU/real-time, ~15 connected UEs before deadline failure) is far below any RAM "
        "ceiling — an AI-pipeline memory reservation would not be the binding constraint for the OAI track.",
    ], Inches(5.2), size=12)
    add_source_label(s, "Sources: OAI phase2 SUMMARY.md/resource_*.csv (REAL MEASUREMENT) + MEMORY_BUDGET_NOTE.md")

    # --- Slide: Final Summary ---
    s = blank_slide(prs)
    add_title(s, "Assignment Status — Baseline Before PAIBO")
    col_w = Inches(4.0)
    headers = ["COMPLETED", "PARTIAL / LIMITATIONS", "NOT YET AVAILABLE"]
    completed = ["Real OAI CU-DU baseline measurements (8-UE pilot + 2-UE registration case)",
                 "ns-3/5G-LENA baseline (7-level ladder, 1-200 UEs)", "6-class traffic model",
                 "CU-DU topology in ns-3", "UE scaling ladder", "KPI extraction", "KPI validation",
                 "Excel/CSV outputs", "OAI vs ns-3 KPI availability analysis"]
    partial = ["ns-3 F1 link is topological, not functional 3GPP F1AP", "F1 heartbeat does not carry actual UE bearer traffic",
               "ns-3 core is LTE-EPC-style, not real 5GC", "Many OAI radio KPIs unavailable — 0/8 UEs attached in the measured pilot",
               "ns-3 traffic uses OnOffApplication, not real iperf3", "Only the IDEAL channel currently validated in ns-3",
               "No 16-UE OAI measured dataset exists — peak real OAI registration is 2 UEs"]
    missing = ["PAIBO Result Type 1 — Bearer Setup Latency", "PAIBO Result Type 2 — DRB/RL-SDAP Reduction",
               "PAIBO Result Type 3 — MAC-CE Adaptation Latency", "PAIBO Result Type 4 — BIP ML Accuracy",
               "Actual PAIBO-vs-baseline experimental comparison"]
    cols_data = [completed, partial, missing]
    top = Inches(1.4)
    for i, (h, items) in enumerate(zip(headers, cols_data)):
        left = Inches(0.4) + i * Inches(4.2)
        hb = s.shapes.add_shape(1, left, top, Inches(4.0), Inches(0.4))
        hb.fill.solid()
        hb.fill.fore_color.rgb = DARK_GRAY if i == 0 else (RGBColor(0xB0, 0x83, 0x00) if i == 1 else RED_FLAG)
        hb.line.fill.background()
        htf = hb.text_frame
        hp = htf.paragraphs[0]
        hp.text = h
        hp.font.bold = True
        hp.font.color.rgb = WHITE
        hp.font.size = Pt(13)
        hp.alignment = PP_ALIGN.CENTER
        box = s.shapes.add_textbox(left, top + Inches(0.5), Inches(4.0), Inches(5.5))
        tf = box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"- {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
    add_source_label(s, "This slide reflects the actual, verified state of this machine's data as of this presentation — nothing extrapolated.",
                      top=Inches(7.1))

    out_path = os.path.join(BASE, "PAIBO_baseline_KPI_comparison.pptx")
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = main()
    prs = Presentation(path)
    n_slides = len(prs.slides)
    print("=== PPTX GENERATION COMPLETE ===")
    print()
    print("File:")
    print(f"  {path}")
    print()
    print("Slides:")
    print(f"  {n_slides}")
    print()
    print("Data sources:")
    print("  OAI CU-DU logs (phase2/20260902_vrtsim_cudu_8ue_106prb/, phase2/20260903_100ue_vrtsim_cudu_189prb_patched/SUMMARY.md)")
    print("  5G-LENA/ns-3 (ns3_cudu_phase/, validated ladder)")
    print("  PAIBO_Patent_Vartika.pptx (reference figures only, clearly labeled, not measurements)")
    print()
    print("Real OAI data:")
    print("  YES")
    print()
    print("ns-3 simulation data:")
    print("  YES")
    print()
    print("PAIBO Result Types 1-4:")
    print("  NOT AVAILABLE")
    print()
    print("BIP accuracy:")
    print("  NOT AVAILABLE")
    print()
    print("No OAI source modified.")
    print("No ns-3 simulation run.")
    print("No GitHub push.")
